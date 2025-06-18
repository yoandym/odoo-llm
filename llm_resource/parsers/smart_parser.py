"""Smart parser implementation using multiple specialized libraries."""

import base64
import io
import logging
import re

from .base_parser import BaseDocumentParser


# Check library availability without importing heavy components at module level
def _check_library_available(library_name):
    """Check if a library is available without importing it at module level."""
    try:
        __import__(library_name)
        return True
    except ImportError:
        return False

# Library availability flags (checked once at module load)
PYMUPDF_AVAILABLE = _check_library_available('fitz')
PANDAS_AVAILABLE = _check_library_available('pandas')
DOCX_AVAILABLE = _check_library_available('docx')
PPTX_AVAILABLE = _check_library_available('pptx')
BS4_AVAILABLE = _check_library_available('bs4')
CAMELOT_AVAILABLE = _check_library_available('camelot')

_logger = logging.getLogger(__name__)


class SmartParser(BaseDocumentParser):
    """Advanced document parser using specialized libraries for different formats."""
    
    @property
    def name(self) -> str:
        return "Smart Parser"
    
    @property
    def description(self) -> str:
        return (
            "Advanced document parser using specialized libraries. "
            "Extracts tables and images from various document formats. "
            "Moderate memory requirements. "
            "Best for most document types including PDF, Word, PowerPoint, and HTML. "
            "Uses PyMuPDF, python-docx, python-pptx and BeautifulSoup."
        )
    
    @property
    def requirements(self) -> str:
        return "Moderate system requirements. Requires multiple parsing libraries."
    
    @property
    def use_cases(self) -> str:
        return (
            "Best for complex documents with tables, images, and structured content. "
            "Supports PDF, DOCX, PPTX, and HTML formats with advanced parsing."
        )
    
    def parse(self, record, field) -> bool:
        """Parse content using specialized libraries based on format."""
        try:
            # Get record name or default to model name and ID
            record_name = (
                record.display_name
                if hasattr(record, "display_name")
                else f"{record._name} #{record.id}"
            )
            
            # Initialize the document content
            content_parts = [f"# {record_name}"]
            
            mimetype = field.get("mimetype", "")
            raw_content = field.get("rawcontent", "")
            
            # Convert to bytes if needed for binary formats
            if isinstance(raw_content, str) and mimetype in [
                "application/pdf",
                "application/octet-stream",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ]:
                try:
                    raw_content = base64.b64decode(raw_content)
                except Exception as decode_error:
                    self._log_error(f"Error decoding binary content: {str(decode_error)}")
                    raw_content = None
            
            # Process based on mimetype
            if mimetype == "application/pdf" and PYMUPDF_AVAILABLE and raw_content:
                content_parts.extend(self._smart_parse_pdf(raw_content, record))
                
            elif (mimetype.endswith(".docx") or "wordprocessingml" in mimetype) and DOCX_AVAILABLE and raw_content:
                content_parts.extend(self._smart_parse_docx(raw_content, record))
                
            elif mimetype.endswith(".pptx") or "presentationml" in mimetype and PPTX_AVAILABLE and raw_content:
                content_parts.extend(self._smart_parse_pptx(raw_content, record))
                
            elif "html" in mimetype and BS4_AVAILABLE:
                # For HTML content
                content_parts.extend(self._smart_parse_html(raw_content, record))
                
            elif mimetype.startswith("text/"):
                # For plain text content
                content_parts.append(raw_content)
                
            else:
                # Default handling for unsupported types
                self._log_error(f"Smart Parser: Unsupported mimetype {mimetype}. Using default fallback.")
                return self._parse_default(record, field)
            
            # Update resource with extracted content
            record.content = "\n\n".join([part for part in content_parts if part])
            return True
            
        except Exception as e:
            self._log_error("Smart Parser error", e)
            return False
    
    def _parse_default(self, record, field):
        """Default parser for unsupported types."""
        try:
            mimetype = field["mimetype"]
            record_name = getattr(record, 'display_name', f"{record._name} #{record.id}")
            record.content = f"""# {record_name}

**File Type**: {mimetype}
**Description**: This file is of type {mimetype} which cannot be directly parsed into text content.
**Source Record**: {record._name} #{record.id}
"""
            return True
        except Exception as e:
            self._log_error("Error in default parsing", e)
            return False

    def _smart_parse_pdf(self, pdf_data, record):
        """Parse PDF content using PyMuPDF with enhanced resume detection and semantic sectioning for better chunking"""
        # Import fitz only when needed
        try:
            import fitz  # PyMuPDF
        except ImportError:
            self._log_error("PyMuPDF (fitz) library not available for PDF parsing")
            return []
            
        content_parts = []

        # Create a temporary file object for PDF data
        pdf_file = io.BytesIO(pdf_data)
        
        # Open PDF with PyMuPDF
        with fitz.open(stream=pdf_file, filetype="pdf") as doc:
            page_count = len(doc)
            
            # Start with main document title
            title = doc.metadata.get('title') if doc.metadata else None
            if title:
                content_parts.append(f"# {title}")
            else:
                record_name = getattr(record, 'display_name', f"{record._name} #{record.id}")
                content_parts.append(f"# {record_name}")
            
            # Add document metadata in its own section for better chunking
            content_parts.append("## Document Information")
            content_parts.append(f"- **Pages**: {page_count}")
            
            # Extract document metadata if available
            metadata = doc.metadata
            if metadata:
                if metadata.get('author'):
                    content_parts.append(f"- **Author**: {metadata.get('author')}")
                if metadata.get('subject'):
                    content_parts.append(f"- **Subject**: {metadata.get('subject')}")
                if metadata.get('keywords'):
                    content_parts.append(f"- **Keywords**: {metadata.get('keywords')}")
                if metadata.get('creationDate'):
                    content_parts.append(f"- **Creation Date**: {metadata.get('creationDate')}")
            
            # Keep track of page and document structure for better chunking
            document_structure = []
            content_blocks = []
            image_count = 0
            
            # First pass: analyze document structure based on formatting and position
            for page_num in range(page_count):
                page = doc[page_num]
                
                # Start a new page in our document structure
                document_structure.append({
                    "page": page_num + 1,
                    "blocks": []
                })
                
                # Extract text blocks with structure recognition
                blocks = page.get_text("blocks")
                
                if not blocks:
                    continue
                    
                # First get page dimensions for position analysis
                page_width = page.rect.width
                page_height = page.rect.height
                
                # Organize blocks by vertical position (y0)
                blocks_by_position = sorted(blocks, key=lambda b: (b[1], b[0]))  # Sort by y0, then x0
                
                # Analyze each block for structural cues
                for block_idx, block in enumerate(blocks_by_position):
                    # Unpack block coordinates and text
                    x0, y0, x1, y1, text, block_type, block_no = block
                    
                    # Skip empty blocks
                    if not text.strip():
                        continue
                    
                    # Calculate some properties for structure detection
                    is_at_top = y0 < page_height * 0.2  # Top 20% of page
                    is_centered = abs((x0 + x1) / 2 - page_width / 2) < page_width * 0.2  # Within 20% of center
                    font_size_large = False
                    
                    # Get more block details if possible
                    try:
                        # Try to get spans which contain font information
                        spans = page.get_text("dict", clip=(x0, y0, x1, y1))
                        if "blocks" in spans and spans["blocks"]:
                            for span_block in spans["blocks"]:
                                if "lines" in span_block:
                                    for line in span_block["lines"]:
                                        if "spans" in line:
                                            for span in line["spans"]:
                                                # Check if font size is larger (relative to other blocks)
                                                if "size" in span and span["size"] > 11:  # Arbitrary threshold
                                                    font_size_large = True
                                                    break
                    except Exception:
                        # If we can't get font information, continue without it
                        pass
                    
                    # Detect potential heading based on formatting
                    is_heading = (
                        # ALL CAPS, relatively short text
                        (text.strip().isupper() and len(text.strip().split()) <= 5) or
                        # Title Case followed by colon is often a heading
                        bool(re.match(r'^[A-Z][a-z]+(\s+[A-Z][a-z]+){0,2}:$', text.strip())) or 
                        # Short text at top of page is likely a header
                        (is_at_top and len(text.strip()) < 80) or
                        # Centered text with larger font is likely a header
                        (is_centered and font_size_large) or
                        # Standalone number or short text that's centered
                        (is_centered and len(text.strip()) < 30) or
                        # Short text followed by a blank line often indicates a header
                        (len(text.strip().split('\n')) == 1 and len(text.strip()) < 50 and font_size_large)
                    )
                    
                    # Determine the semantic structure
                    block_role = "heading" if is_heading else "content"
                    
                    # Add to document structure
                    document_structure[-1]["blocks"].append({
                        "text": text.strip(),
                        "role": block_role,
                        "pos": (x0, y0, x1, y1),
                        "is_centered": is_centered,
                        "is_at_top": is_at_top,
                        "has_large_font": font_size_large
                    })
                    
                    # Store for content generation
                    content_blocks.append({
                        "page": page_num + 1,
                        "text": text.strip(),
                        "role": block_role
                    })
            
            # Now generate content with proper markdown structure for good chunking
                # Process each page and create semantic sections
                content_parts.append("## Document Content")
                
                # Generate page-specific content with proper semantic structure
                current_page = None
                heading_level = 3  # Start with ### for main section headings
                
                # Process blocks sequentially for content
                for block in content_blocks:
                    page_num = block["page"]
                    text = block["text"]
                    role = block["role"]
                    
                    # Add page break if we're moving to a new page
                    if page_num != current_page:
                        content_parts.append(f"\n### Page {page_num}")
                        current_page = page_num
                    
                    # Format based on the block's role
                    if role == "heading":
                        # Use proper heading level to create semantic sections
                        # This creates breakpoints for the chunker
                        content_parts.append(f"\n{'#' * heading_level} {text}")
                    else:
                        # Process content blocks
                        lines = text.split('\n')
                        formatted_lines = []
                        in_bullet_list = False
                        
                        for line in lines:
                            line = line.strip()
                            if not line:
                                if formatted_lines and formatted_lines[-1]:
                                    formatted_lines.append("")
                                in_bullet_list = False
                                continue
                            
                            # Check for bullet points with a more general pattern
                            bullet_match = re.match(r'^[•\-\*∙⦁◦○●◆■□▪▫]|^\d+[\.\)]|^[a-z][\.\)]|^[\(]?[ivxIVX]+[\.\)]', line)
                            
                            if bullet_match:
                                if not in_bullet_list:
                                    formatted_lines.append("")  # Add spacing before list
                                in_bullet_list = True
                                # Ensure bullet points are formatted consistently
                                formatted_lines.append(f"- {line[bullet_match.end():].strip()}")
                            else:
                                # Regular paragraph text
                                if in_bullet_list:
                                    # If it's an indented continuation of a bullet point
                                    if (line.startswith("  ") or line.startswith("\t")) and formatted_lines:
                                        formatted_lines[-1] += f" {line.strip()}"
                                    else:
                                        in_bullet_list = False
                                        formatted_lines.append(line)
                                else:
                                    formatted_lines.append(line)
                        
                        # Add the formatted content to our parts
                        content_parts.append("\n".join(formatted_lines))
                
                # Now process each page for tables and images
                for page_num in range(page_count):
                    page = doc[page_num]
                    
                    # Extract tables using camelot if available
                    if CAMELOT_AVAILABLE:
                        try:
                            import tempfile

                            # Save PDF to temp file for camelot
                            with tempfile.NamedTemporaryFile(suffix='.pdf') as tmp:
                                tmp.write(pdf_data)
                                tmp.flush()
                                
                                # Extract tables for this page
                                tables = camelot.read_pdf(tmp.name, pages=str(page_num + 1))
                                if len(tables) > 0:
                                    # Each table gets its own section for better chunking
                                    content_parts.append(f"\n### Tables - Page {page_num + 1}")
                                    for i, table in enumerate(tables):
                                        # Give each table its own headline for better chunking
                                        content_parts.append(f"\n#### Table {i+1}")
                                        
                                        if PANDAS_AVAILABLE:
                                            df = table.df
                                            if not df.empty:
                                                # Add descriptive text about the table
                                                content_parts.append(f"Table with {len(df)} rows and {len(df.columns)} columns")
                                                
                                                # Format table with proper markdown
                                                table_md = df.to_markdown(index=False)
                                                if table_md:
                                                    content_parts.append(f"\n{table_md}\n")
                                        else:
                                            # Alternative table formatting if pandas isn't available
                                            row_count = len(table.data)
                                            content_parts.append(f"Table with {row_count} rows")
                                            
                                            # Header row
                                            if row_count > 0:
                                                header = "| " + " | ".join(table.data[0]) + " |"
                                                content_parts.append(header)
                                                
                                                # Table separator
                                                separator = "|" + "|".join(["----" for _ in table.data[0]]) + "|"
                                                content_parts.append(separator)
                                                
                                                # Data rows
                                                for row in table.data[1:]:
                                                    content_parts.append("| " + " | ".join(row) + " |")
                        except Exception as table_error:
                            _logger.warning(f"Error extracting tables from page {page_num + 1}: {str(table_error)}")
                
                    # Extract images as a separate section for better chunking
                    try:
                        image_list = page.get_images(full=True)
                        if image_list:
                            # Create a separate section for images to improve chunking
                            content_parts.append(f"\n### Images - Page {page_num + 1}")
                            content_parts.append(f"Found {len(image_list)} images on this page")
                        
                        for img_index, img in enumerate(image_list):
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            
                            if base_image:
                                image_data = base_image["image"]
                                image_ext = base_image["ext"]
                                image_name = f"image_p{page_num + 1}_{img_index + 1}.{image_ext}"
                                
                                # Create attachment for the image
                                img_attachment = record.env["ir.attachment"].create({
                                    "name": image_name,
                                    "datas": base64.b64encode(image_data),
                                    "res_model": "llm.resource",
                                    "res_id": record.id,
                                    "mimetype": f"image/{image_ext}",
                                })
                                
                                # Add image reference to markdown content - each as a separate item
                                if img_attachment:
                                    image_url = f"/web/image/{img_attachment.id}"
                                    # Give each image its own heading for better chunking
                                    content_parts.append(f"\n#### Image {img_index + 1}")
                                    content_parts.append(f"![{image_name}]({image_url})")
                                    image_count += 1
                    except Exception as img_error:
                        _logger.warning(f"Error extracting images from page {page_num + 1}: {str(img_error)}")
            
            # Create a dedicated document summary section for better chunking
            content_parts.append("\n## Document Summary")
            content_parts.append(f"- **Document name**: {record.display_name}")
            content_parts.append(f"- **Total pages**: {page_count}")
            if image_count > 0:
                content_parts.append(f"- **Total images**: {image_count}")
            if metadata and metadata.get('title'):
                content_parts.append(f"- **Document title**: {metadata.get('title')}")
                
            # Add document structure information to help with chunking
            content_parts.append("\n### Document Structure")
            
            # Summarize the document structure (number of headings, content sections, etc.)
            heading_count = sum(1 for block in content_blocks if block["role"] == "heading")
            content_parts.append(f"- **Headings**: {heading_count}")
            content_parts.append(f"- **Content blocks**: {len(content_blocks) - heading_count}")
            content_parts.append(f"- **Tables and images**: {image_count}")
        
        return content_parts

    def _smart_parse_docx(self, docx_data, record):
        """Parse DOCX content using python-docx"""
        # Import docx only when needed
        try:
            import docx
        except ImportError:
            self._log_error("python-docx library not available for DOCX parsing")
            return []
        
        # Import pandas if available for table formatting
        pd = None
        if PANDAS_AVAILABLE:
            try:
                import pandas as pd
            except ImportError:
                pd = None
            
        content_parts = []

        # Create a temporary file object
        docx_file = io.BytesIO(docx_data)
        
        # Open the document
        doc = docx.Document(docx_file)
        
        # Document title
        record_name = getattr(record, 'display_name', f"{record._name} #{record.id}")
        content_parts.append(f"# {record_name}")
        
        # Extract document properties in its own section
        content_parts.append("\n## Document Information")
        
        # Add core properties
        core_props = doc.core_properties
        if core_props.title:
            content_parts.append(f"- **Title**: {core_props.title}")
        
        if core_props.author:
            content_parts.append(f"- **Author**: {core_props.author}")
            
        if core_props.created:
            content_parts.append(f"- **Created**: {core_props.created}")
            
        if core_props.modified:
            content_parts.append(f"- **Modified**: {core_props.modified}")
            
        if core_props.subject:
            content_parts.append(f"- **Subject**: {core_props.subject}")
        
        # Track current section for better organization
        current_section = None
        
        # Process document content, organizing by sections
        content_parts.append("\n## Document Content")
        
        # Process paragraphs
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
                
            # Check if paragraph is a heading
            if para.style and para.style.name and para.style.name.startswith('Heading'):
                heading_level = int(para.style.name.replace('Heading ', ''))
                # Add additional # for proper nesting (h1 = ###, h2 = ####, etc.)
                content_parts.append(f"{'#' * (heading_level + 2)} {para.text}")
                current_section = para.text
            else:
                # Detect if this might be an unlabeled section header
                if (para.text.isupper() and len(para.text.split()) <= 5) or re.match(r'^[A-Z][a-z]+(\s+[A-Z][a-z]+){0,2}:$', para.text):
                    content_parts.append(f"\n### {para.text}")
                    current_section = para.text
                else:
                    # For regular paragraphs, see if we can identify bullet points
                    text = para.text.strip()
                    bullet_match = re.match(r'^[•\-\*∙⦁◦]|^\d+[\.\)]|^[a-z][\.\)]', text)
                    
                    if bullet_match:
                        # Format as a bullet point
                        content_parts.append(f"- {text[bullet_match.end():].strip()}")
                    else:
                        # Regular paragraph
                        content_parts.append(text)
        
        # Process tables in a dedicated section
        if doc.tables:
            content_parts.append("\n## Tables")
            for i, table in enumerate(doc.tables):
                content_parts.append(f"\n### Table {i+1}")
                
                table_data = []
                for j, row in enumerate(table.rows):
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    # Add table metadata
                    content_parts.append(f"Table with {len(table_data)} rows and {len(table_data[0]) if table_data else 0} columns")
                    
                    if pd is not None:
                        # Using pandas for better table formatting
                        try:
                            df = pd.DataFrame(table_data[1:], columns=table_data[0] if table_data else None)
                            content_parts.append(f"\n{df.to_markdown(index=False)}\n")
                        except Exception:
                            # Fallback to manual formatting if pandas fails
                            self._manual_table_format(content_parts, table_data)
                    else:
                        # Manual table formatting
                        self._manual_table_format(content_parts, table_data)
        
        # Extract images
        image_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    image_name = f"docx_image_{image_count + 1}.png"
                    
                    # Create attachment for the image
                    img_attachment = record.env["ir.attachment"].create({
                        "name": image_name,
                        "datas": base64.b64encode(image_data),
                        "res_model": "llm.resource",
                        "res_id": record.id,
                        "mimetype": "image/png",
                    })
                    
                    # Add image reference to markdown content
                    if img_attachment:
                        image_url = f"/web/image/{img_attachment.id}"
                        content_parts.append(f"\n![{image_name}]({image_url})\n")
                        image_count += 1
                except Exception as img_error:
                    _logger.warning(f"Error extracting DOCX image {image_count}: {str(img_error)}")
        
        if image_count > 0:
            content_parts.append(f"\n## Document contains {image_count} images\n")
        
        return content_parts

    def _manual_table_format(self, content_parts, table_data):
        """Format table data manually without pandas"""
        if table_data:
            # Header row
            content_parts.append("| " + " | ".join(table_data[0]) + " |")
            # Table separator
            content_parts.append("|" + "|".join(["----" for _ in table_data[0]]) + "|")
            # Data rows
            for row in table_data[1:]:
                content_parts.append("| " + " | ".join(row) + " |")

    def _smart_parse_pptx(self, pptx_data, record):
        """Parse PPTX content using python-pptx with improved structure for chunking"""
        # Import pptx only when needed
        try:
            from pptx import Presentation
        except ImportError:
            self._log_error("python-pptx library not available for PPTX parsing")
            return []
            
        content_parts = []

        # Create a temporary file object
        pptx_file = io.BytesIO(pptx_data)
        
        # Open the presentation - use imported Presentation class
        prs = Presentation(pptx_file)
        
        # Document title
        record_name = getattr(record, 'display_name', f"{record._name} #{record.id}")
        content_parts.append(f"# {record_name}")
        
        # Add presentation metadata
        content_parts.append("\n## Presentation Information")
        content_parts.append("- **Type**: PowerPoint Presentation")
        content_parts.append(f"- **Slides**: {len(prs.slides)}")
        
        # Try to extract properties if available
        try:
            if hasattr(prs.core_properties, 'title') and prs.core_properties.title:
                content_parts.append(f"- **Title**: {prs.core_properties.title}")
            if hasattr(prs.core_properties, 'author') and prs.core_properties.author:
                content_parts.append(f"- **Author**: {prs.core_properties.author}")
            if hasattr(prs.core_properties, 'subject') and prs.core_properties.subject:
                content_parts.append(f"- **Subject**: {prs.core_properties.subject}")
        except Exception:
            # Continue even if properties extraction fails
            pass
            
        # Process slides - each slide gets its own section for better chunking
        content_parts.append("\n## Presentation Content")
        
        total_images = 0
        
        # Process slides
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            
            # Create a dedicated section for each slide
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text
                content_parts.append(f"\n### Slide {slide_num}: {title}")
            else:
                content_parts.append(f"\n### Slide {slide_num}")
            
            # Extract text from all shapes
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:  # Skip title as we already added it
                        texts.append(shape.text)
            
            if texts:
                content_parts.append("#### Slide Content")
                # Add each text element with proper formatting
                for j, text in enumerate(texts):
                    # Check if text might be a bullet point
                    lines = text.split('\n')
                    formatted_lines = []
                    
                    for line in lines:
                        line = line.strip()
                        if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                            # Format as markdown bullet
                            formatted_lines.append(f"- {line[1:].strip()}")
                        else:
                            formatted_lines.append(line)
                    
                    content_parts.append("\n".join(formatted_lines))
            
            # Extract images in a dedicated subsection
            image_count_slide = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 13 is the value for picture
                    try:
                        # If first image on slide, add header
                        if image_count_slide == 0:
                            content_parts.append("\n#### Slide Images")
                        
                        image = shape.image
                        image_bytes = image.blob
                        image_name = f"slide{slide_num}_image_{image_count_slide + 1}.png"
                        
                        # Create attachment for the image
                        img_attachment = record.env["ir.attachment"].create({
                            "name": image_name,
                            "datas": base64.b64encode(image_bytes),
                            "res_model": "llm.resource",
                            "res_id": record.id,
                            "mimetype": "image/png",
                        })
                        
                        # Add image reference to markdown content - each gets its own subsection
                        if img_attachment:
                            image_url = f"/web/image/{img_attachment.id}"
                            content_parts.append(f"##### Image {image_count_slide + 1}")
                            content_parts.append(f"![{image_name}]({image_url})")
                            image_count_slide += 1
                            total_images += 1
                    except Exception as img_error:
                        _logger.warning(f"Error extracting PPTX image on slide {slide_num}: {str(img_error)}")
            
            if image_count_slide > 0:
                content_parts.append(f"\n*Slide contains {image_count_slide} images*")
        
        # Add presentation summary
        if total_images > 0:
            content_parts.append("\n## Presentation Summary")
            content_parts.append(f"- **Total slides**: {len(prs.slides)}")
            content_parts.append(f"- **Total images**: {total_images}")
            
        return content_parts

    def _smart_parse_html(self, html_content, record):
        """Parse HTML content using BeautifulSoup with enhanced semantic structure"""
        # Import BeautifulSoup only when needed
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            self._log_error("BeautifulSoup4 library not available for HTML parsing")
            return []
        
        # Import pandas if available for table formatting
        pd = None
        if PANDAS_AVAILABLE:
            try:
                import pandas as pd
            except ImportError:
                pd = None
            
        content_parts = []
        from urllib.parse import urljoin, urlparse

        # Parse HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Document title as main heading
        record_name = getattr(record, 'display_name', f"{record._name} #{record.id}")
        if soup.title and soup.title.string:
            content_parts.append(f"# {soup.title.string.strip()}")
        else:
            content_parts.append(f"# {record_name}")
        
        # Extract metadata in a dedicated section
        content_parts.append("\n## Document Information")
        content_parts.append("- **Source**: HTML Document")
        
        # If there's meta description, add it
        meta_desc = soup.find('meta', attrs={'name': 'description'})
        if meta_desc and meta_desc.get('content'):
            content_parts.append(f"- **Description**: {meta_desc.get('content')}")
            
        # Extract any other useful metadata
        meta_keywords = soup.find('meta', attrs={'name': 'keywords'})
        if meta_keywords and meta_keywords.get('content'):
            content_parts.append(f"- **Keywords**: {meta_keywords.get('content')}")
            
        # Extract main content
        body = soup.body if soup.body else soup
        
        # Process main content
        content_parts.append("\n## Document Content")
        
        # First collect all headings to understand document structure
        headings = body.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
        
        # If document has headings, process by sections
        if headings:
            for heading in headings:
                # Get heading level and add 2 to nest properly (h1 -> ###, h2 -> ####, etc.)
                heading_level = min(int(heading.name[1]) + 2, 6)  # Max 6 levels of nesting
                heading_text = heading.get_text().strip()
                content_parts.append(f"\n{'#' * heading_level} {heading_text}")
                
                # Collect content until next heading
                content_elements = []
                current = heading.next_sibling
                
                # Check if we've reached another heading
                while current and not (hasattr(current, 'name') and current.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                    if hasattr(current, 'name') and current.name:
                        # Process different HTML elements
                        if current.name == 'p':
                            text = current.get_text().strip()
                            if text:
                                content_elements.append(text)
                        elif current.name in ['ul', 'ol']:
                            # Process lists
                            for li in current.find_all('li'):
                                li_text = li.get_text().strip()
                                if li_text:
                                    content_elements.append(f"- {li_text}")
                        elif current.name == 'img' and current.get('src'):
                            try:
                                # Handle image
                                img_url = current.get('src')
                                
                                # Process relative URLs
                                if not bool(urlparse(img_url).netloc):
                                    # Find base URL
                                    base_tag = soup.find('base')
                                    if base_tag and base_tag.get('href'):
                                        base_url = base_tag.get('href')
                                        img_url = urljoin(str(base_url), img_url)
                                    else:
                                        # Can't resolve relative URL
                                        content_elements.append(f"*Image: {img_url} (relative path)*")
                                        current = current.next_sibling
                                        continue
                                
                                # Add image reference
                                alt_text = current.get('alt', 'Image')
                                content_elements.append(f"\n#### Image: {alt_text}")
                                content_elements.append(f"![{alt_text}]({img_url})")
                                
                            except Exception as img_error:
                                _logger.warning(f"Error processing HTML image: {str(img_error)}")
                                content_elements.append(f"*Image reference: {img_url}*")
                        
                        elif current.name == 'table':
                            try:
                                # Process table into dedicated section
                                content_elements.append("\n#### Table")
                                
                                # Extract table data
                                rows = []
                                for tr in current.find_all('tr'):
                                    row = []
                                    for td in tr.find_all(['td', 'th']):
                                        row.append(td.get_text().strip())
                                    if row:
                                        rows.append(row)
                                
                                # Format table in markdown
                                if rows:
                                    content_elements.append(f"Table with {len(rows)} rows")
                                    
                                    if pd is not None:
                                        # Use pandas for formatted tables
                                        try:
                                            has_headers = bool(current.find('th'))
                                            if has_headers and len(rows) > 1:
                                                df = pd.DataFrame(rows[1:], columns=rows[0])
                                                content_elements.append(f"\n{df.to_markdown(index=False)}\n")
                                            else:
                                                df = pd.DataFrame(rows)
                                                content_elements.append(f"\n{df.to_markdown(index=False)}\n")
                                        except Exception:
                                            # Fallback to manual formatting if pandas fails
                                            self._manual_table_format(content_elements, rows)
                                    else:
                                        # Manual table formatting
                                        self._manual_table_format(content_elements, rows)
                            except Exception as table_error:
                                _logger.warning(f"Error processing HTML table: {str(table_error)}")
                                content_elements.append("*Table could not be processed*")
                        
                        # Add more element types as needed
                        
                    current = current.next_sibling
                
                # Add collected content if any
                if content_elements:
                    content_parts.append("\n".join(content_elements))
        else:
            # No headings found, process paragraphs directly
            paragraphs = []
            for para in body.find_all('p'):
                text = para.get_text().strip()
                if text:
                    paragraphs.append(text)
            
            if paragraphs:
                content_parts.append("\n".join(paragraphs))
            else:
                # Fall back to just getting text from body
                body_text = body.get_text().strip()
                if body_text:
                    content_parts.append(body_text)
        
        return content_parts
