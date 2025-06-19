"""Smart parser implementation using multiple specialized libraries."""

import base64
import io
import json
import logging
import re

from odoo import _, exceptions

from .base_parser import BaseDocumentParser


# Check library availability without importing heavy components at module level
def _check_library_available(library_name):
    """Check if a library is available without importing it at module level."""
    try:
        __import__(library_name)
        return True
    except ImportError:
        return False


# Library availability flags (checked once at module load)
PYMUPDF_AVAILABLE = _check_library_available('fitz')
PANDAS_AVAILABLE = _check_library_available('pandas')
DOCX_AVAILABLE = _check_library_available('docx')
PPTX_AVAILABLE = _check_library_available('pptx')
CAMELOT_AVAILABLE = _check_library_available('camelot')
MD_AVAILABLE = _check_library_available('markdownify')

_logger = logging.getLogger(__name__)


class DefaultParser(BaseDocumentParser):
    """Default Document parser using specialized libraries for different formats."""
    
    @property
    def name(self) -> str:
        return "Default Parser"
    
    @property
    def description(self) -> str:
        return (
            "Default document parser using specialized libraries. "
            "Extracts tables and images from various document formats. "
            "Moderate memory requirements. "
            "Best for most document types including PDF, Word, PowerPoint, Markdown and HTML. "
            "Uses PyMuPDF, python-docx, python-pptx, markdownify and BeautifulSoup."
        )
    
    @property
    def requirements(self) -> str:
        return "Moderate system requirements. Requires multiple parsing libraries."
    
    @property
    def use_cases(self) -> str:
        return (
            "Best for complex documents with tables, images, and structured content. "
            "Supports PDF, DOCX, PPTX, Markdown and HTML formats with advanced parsing."
        )
    
    def _get_content_type_handler(self, mimetype, filename=""):
        """Determine the appropriate handler for a given content type"""
        
        # Check filename for extensions
        is_markdown = filename and filename.lower().endswith(('.md', '.markdown'))
        
        if mimetype == "application/pdf" and PYMUPDF_AVAILABLE:
            return self._smart_parse_pdf
            
        elif (mimetype.endswith(".docx") or "wordprocessingml" in mimetype) and DOCX_AVAILABLE:
            return self._smart_parse_docx
            
        elif mimetype.endswith(".pptx") or "presentationml" in mimetype and PPTX_AVAILABLE:
            return self._smart_parse_pptx
            
        elif "html" in mimetype and MD_AVAILABLE:
            return self._parse_html
            
        # special case for markdown files sometimes detected as octet-stream
        elif (mimetype == "application/octet-stream" and is_markdown) or mimetype == "text/markdown":
            return self._parse_markdown
            
        elif mimetype.startswith("text/"):
            return self._parse_text
            
        elif mimetype.startswith("image/"):
            return self._parse_image
            
        elif mimetype == "application/json":
            return self._parse_json
            
        else:
            return self._parse_default
    
    def parse(self, resource, field) -> bool:
        """Parse content using specialized libraries based on format."""
        try:
            # Get field information
            field_name = field.get("field_name", "")
            mimetype = field.get("mimetype", "")
            raw_content = field.get("rawcontent", "")
            
            # Skip empty content
            if not raw_content:
                _logger.warning(f"Skipping empty field: {field_name}")
                return False
            
            # Get record name for context
            record_name = resource.name or f"Resource #{resource.id}"
            
            # Get the appropriate handler for this content type
            handler = self._get_content_type_handler(mimetype, filename=field_name)
            
            # Process the field with the appropriate handler
            processed_content = handler(resource, field)
            
            # Update the resource content with the processed content
            # Note: handlers are expected to return the processed content as string
            if processed_content:
                # If content already exists, append to it
                if resource.content:
                    resource.content = f"{resource.content}\n\n{processed_content}"
                else:
                    resource.content = processed_content
                return True
                
            return False
                
        except Exception as e:
            self._log_error(f"Error parsing field {field.get('field_name')}", exception=e)
            # This would normally use resource._post_styled_message but that's handled by caller
            return False

    def _smart_parse_pdf(self, resource, field):
        """Parse PDF content using PyMuPDF with enhanced resume detection and semantic sectioning."""
        # Import fitz only when needed
        try:
            import fitz  # PyMuPDF
        except ImportError:
            self._log_error("PyMuPDF (fitz) library not available for PDF parsing")
            return None
            
        content_parts = []
        field_name = field.get("field_name", "")
        
        # Add field name as header
        content_parts.append(f"## {field_name}")

        # Convert to bytes if needed for binary formats
        raw_content = field.get("rawcontent", "")
        if isinstance(raw_content, str):
            try:
                raw_content = base64.b64decode(raw_content)
            except Exception as e:
                self._log_error(f"Failed to decode PDF content: {str(e)}")
                return f"## {field_name}\n\nError: Could not decode PDF content"
        
        # Create a temporary file object for PDF data
        pdf_file = io.BytesIO(raw_content)
        
        try:
            # Open PDF with PyMuPDF
            with fitz.open(stream=pdf_file, filetype="pdf") as doc:
                page_count = len(doc)
                
                # Get document metadata
                metadata = doc.metadata
                if metadata:
                    content_parts.append("### Document Information")
                    content_parts.append(f"- **Pages**: {page_count}")
                    if metadata.get('title'):
                        content_parts.append(f"- **Title**: {metadata.get('title')}")
                    if metadata.get('author'):
                        content_parts.append(f"- **Author**: {metadata.get('author')}")
                    if metadata.get('subject'):
                        content_parts.append(f"- **Subject**: {metadata.get('subject')}")
                    if metadata.get('keywords'):
                        content_parts.append(f"- **Keywords**: {metadata.get('keywords')}")
                    if metadata.get('creationDate'):
                        content_parts.append(f"- **Creation Date**: {metadata.get('creationDate')}")
                
                # Extract content from each page
                for page_num in range(page_count):
                    page = doc[page_num]
                    
                    # Extract text
                    text = page.get_text()
                    if text.strip():
                        content_parts.append(f"### Page {page_num + 1}")
                        content_parts.append(text.strip())
                    
                    # Extract images
                    # Note: Image extraction is resource-intensive, so limit it for large documents
                    if page_count <= 50:  # Only extract images for smaller documents
                        image_list = page.get_images(full=True)
                        for img_index, img in enumerate(image_list):
                            xref = img[0]
                            try:
                                base_image = doc.extract_image(xref)
                                if base_image:
                                    # Store image as attachment
                                    image_data = base_image["image"]
                                    image_ext = base_image["ext"]
                                    image_name = f"image_{page_num}_{img_index}.{image_ext}"

                                    # Create attachment for the image
                                    img_attachment = resource.env["ir.attachment"].create(
                                        {
                                            "name": image_name,
                                            "datas": base64.b64encode(image_data),
                                            "res_model": "llm.resource",
                                            "res_id": resource.id,
                                            "mimetype": f"image/{image_ext}",
                                        }
                                    )

                                    # Add image reference to markdown content
                                    if img_attachment:
                                        image_url = f"/web/image/{img_attachment.id}"
                                        content_parts.append(f"\n![{image_name}]({image_url})\n")
                            except Exception as e:
                                self._log_error(f"Error extracting image: {str(e)}")
        
        except Exception as e:
            self._log_error(f"Error parsing PDF: {str(e)}")
            return f"## {field_name}\n\nError: Could not parse PDF content: {str(e)}"
            
        # Join all content parts into a single string
        return "\n\n".join(content_parts)

    def _smart_parse_docx(self, resource, field):
        """Parse DOCX files using python-docx."""
        try:
            import docx
        except ImportError:
            self._log_error("python-docx library not available for DOCX parsing")
            return None
        
        field_name = field.get("field_name", "")
        
        # Add field name as header
        content_parts = [f"## {field_name}"]

        # Convert to bytes if needed for binary formats
        raw_content = field.get("rawcontent", "")
        if isinstance(raw_content, str):
            try:
                raw_content = base64.b64decode(raw_content)
            except Exception as e:
                self._log_error(f"Failed to decode DOCX content: {str(e)}")
                return f"## {field_name}\n\nError: Could not decode DOCX content"
        
        # Create a temporary file object
        docx_file = io.BytesIO(raw_content)
        
        try:
            # Open DOCX with python-docx
            doc = docx.Document(docx_file)
            
            # Extract document properties if available
            try:
                core_props = doc.core_properties
                content_parts.append("### Document Information")
                
                if hasattr(core_props, 'title') and core_props.title:
                    content_parts.append(f"- **Title**: {core_props.title}")
                if hasattr(core_props, 'author') and core_props.author:
                    content_parts.append(f"- **Author**: {core_props.author}")
                if hasattr(core_props, 'subject') and core_props.subject:
                    content_parts.append(f"- **Subject**: {core_props.subject}")
                if hasattr(core_props, 'keywords') and core_props.keywords:
                    content_parts.append(f"- **Keywords**: {core_props.keywords}")
                if hasattr(core_props, 'created') and core_props.created:
                    content_parts.append(f"- **Created**: {core_props.created}")
                if hasattr(core_props, 'modified') and core_props.modified:
                    content_parts.append(f"- **Modified**: {core_props.modified}")
            except Exception as e:
                self._log_error(f"Error extracting document properties: {str(e)}")
            
            # Extract document content
            content_parts.append("### Document Content")
            
            # Parse paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check if this is a heading
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.replace('Heading', '')) if para.style.name != 'Heading' else 1
                        # Add the appropriate markdown heading level
                        heading_marks = '#' * (level + 2)  # +2 because we already have ## for the field name
                        content_parts.append(f"{heading_marks} {para.text.strip()}")
                    else:
                        content_parts.append(para.text.strip())
            
            # Parse tables
            for i, table in enumerate(doc.tables):
                content_parts.append(f"#### Table {i+1}")
                
                # Extract table data
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                # Convert to markdown table
                if table_data:
                    # Create header row
                    md_table = ["| " + " | ".join(table_data[0]) + " |"]
                    # Add separator row
                    md_table.append("| " + " | ".join(["---"] * len(table_data[0])) + " |")
                    # Add data rows
                    for row in table_data[1:]:
                        md_table.append("| " + " | ".join(row) + " |")
                    
                    content_parts.append("\n".join(md_table))
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            self._log_error(f"Error parsing DOCX: {str(e)}")
            return f"## {field_name}\n\nError: Could not parse DOCX content: {str(e)}"
    
    def _smart_parse_pptx(self, resource, field):
        """Parse PPTX files using python-pptx."""
        try:
            import pptx
        except ImportError:
            self._log_error("python-pptx library not available for PPTX parsing")
            return None
        
        field_name = field.get("field_name", "")
        
        # Add field name as header
        content_parts = [f"## {field_name}"]

        # Convert to bytes if needed for binary formats
        raw_content = field.get("rawcontent", "")
        if isinstance(raw_content, str):
            try:
                raw_content = base64.b64decode(raw_content)
            except Exception as e:
                self._log_error(f"Failed to decode PPTX content: {str(e)}")
                return f"## {field_name}\n\nError: Could not decode PPTX content"
        
        # Create a temporary file object
        pptx_file = io.BytesIO(raw_content)
        
        try:
            # Open PPTX with python-pptx
            presentation = pptx.Presentation(pptx_file)
            
            # Extract presentation content
            content_parts.append(f"### Presentation Content ({len(presentation.slides)} slides)")
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                content_parts.append(f"#### Slide {i+1}")
                
                # Extract slide title if available
                if slide.shapes.title and slide.shapes.title.text:
                    content_parts.append(f"**Title**: {slide.shapes.title.text}")
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    content_parts.append("\n".join(slide_text))
                else:
                    content_parts.append("*No text content in this slide*")
            
            return "\n\n".join(content_parts)
            
        except Exception as e:
            self._log_error(f"Error parsing PPTX: {str(e)}")
            return f"## {field_name}\n\nError: Could not parse PPTX content: {str(e)}"
    
    def _parse_text(self, resource, field):
        """Parse plain text content."""
        raw_content = field.get("rawcontent", "")
        
        # If content is in bytes, decode it properly
        if isinstance(raw_content, bytes):
            return self._safe_decode(raw_content)
        elif isinstance(raw_content, str) and raw_content.startswith("b'") and raw_content.endswith("'"):
            # This handles the case where bytes are represented as a string like b'content'
            try:
                # Remove the b'' wrapper and parse
                content_str = raw_content[2:-1]
                # Handle escape sequences
                content_str = content_str.encode('utf-8').decode('unicode_escape')
                return content_str
            except Exception as e:
                self._log_error(f"Failed to decode content string representation: {str(e)}")
        
        return raw_content
        
    def _parse_html(self, resource, field):
        """Parse HTML content using markdownify."""
        field_name = field.get("field_name", "")
        raw_content = field.get("rawcontent", "")
        
        # Import markdownify only when needed
        try:
            from markdownify import markdownify as md
            markdown_content = md(raw_content)
            # Format with field name header
            return f"## {field_name}\n\n{markdown_content}"
        except ImportError:
            _logger.warning("Markdownify library not available, returning raw HTML")
            return f"## {field_name}\n\n```html\n{raw_content}\n```"
            
    def _parse_image(self, resource, field):
        """Parse image references."""
        field_name = field.get("field_name", "")
        
        # For attachments that are already in the database
        if hasattr(resource, "attachment_id") and resource.attachment_id:
            image_url = f"/web/image/{resource.attachment_id.id}"
            return f"## {field_name}\n\n![{field_name}]({image_url})"
            
        # For binary fields that are part of the record
        # This would require additional handling to create an attachment
        # if resource.res_model and resource.res_id:
        #    Create attachment logic would go here
            
        return f"## {field_name}\n\n*Image content available but not rendered in markdown*"
        
    def _parse_json(self, resource, field):
        """Parse JSON content."""
        field_name = field.get("field_name", "")
        raw_content = field.get("rawcontent", "")
        
        # Try to parse and pretty print the JSON
        try:
            if isinstance(raw_content, str):
                json_data = json.loads(raw_content)
            else:
                json_data = raw_content
                
            formatted_json = json.dumps(json_data, indent=2, default=str)
            return f"## {field_name}\n\n```json\n{formatted_json}\n```"
        except Exception as e:
            _logger.warning(f"Error parsing JSON content: {e}")
            if isinstance(raw_content, str):
                return f"## {field_name}\n\n{raw_content}"
            return f"## {field_name}\n\n*JSON content could not be parsed*"
            
    def _parse_default(self, resource, field):
        """Default handler for unsupported content types."""
        field_name = field.get("field_name", "")
        mimetype = field.get("mimetype", "")
        
        return f"""## {field_name}

**File Type**: {mimetype}
**Description**: This content is of type {mimetype} which cannot be directly parsed into text.
"""

    def parse_record_fields(self, resource, record):
        """
        Enhanced record field parsing that handles text, binary fields and attachments.
        
        Args:
            resource: The LLM resource record
            record: The Odoo record to extract fields from
            
        Returns:
            List[Dict]: List of field dictionaries with field_name, mimetype, and rawcontent
        """
        results = []
        env = record.env

        # Start with the record name/display_name if available
        record_name_field = (
            "display_name" if hasattr(record, "display_name") else "name"
        )
        record_name = (
            record[record_name_field]
            if hasattr(record, record_name_field)
            else f"{record._name} #{record.id}"
        )
        if record_name:
            results.append(
                {
                    "field_name": record_name_field,
                    "mimetype": "text/plain",
                    "rawcontent": record_name,
                }
            )

        # Process regular fields based on field types
        for field_name, field in record._fields.items():
            try:
                # Skip binary fields (handled separately), computed fields without store, and special fields
                if (
                    field_name.startswith("_") 
                    or field.type == "binary" 
                    or field_name in ["id", record_name_field, "create_uid", "write_uid", "create_date", "write_date"]
                    or (field.compute and not field.store)
                ):
                    continue

                # Skip empty values
                if not record[field_name]:
                    continue

                # Process based on field type
                if field.type == "many2one" and record[field_name]:
                    # For many2one fields, include the name of the related record
                    related_record = record[field_name]
                    if hasattr(related_record, "display_name"):
                        results.append({
                            "field_name": field_name,
                            "mimetype": "text/plain",
                            "rawcontent": f"{field.string}: {related_record.display_name}"
                        })
                
                elif field.type in ["many2many", "one2many"]:
                    # For relational fields, include a list of related record names
                    related_records = record[field_name]
                    if related_records:
                        records_text = "\n".join([f"- {r.display_name}" for r in related_records if hasattr(r, "display_name")])
                        if records_text:
                            results.append({
                                "field_name": field_name,
                                "mimetype": "text/plain",
                                "rawcontent": f"{field.string}:\n{records_text}"
                            })
                
                elif field.type == "html":
                    # Handle HTML fields using markdownify
                    results.append({
                        "field_name": field_name,
                        "mimetype": "text/html",
                        "rawcontent": record[field_name]
                    })
                
                elif field.type == "text":
                    # Handle text fields
                    results.append({
                        "field_name": field_name,
                        "mimetype": "text/plain",
                        "rawcontent": record[field_name]
                    })
                
                else:
                    # Handle other field types
                    results.append({
                        "field_name": field_name,
                        "mimetype": "text/plain",
                        "rawcontent": str(record[field_name])
                    })
            
            except Exception as e:
                _logger.error(f"Error processing field {field_name}: {e}")
                continue

        # Handle binary fields separately
        binary_fields = [f.name for f in record._fields.values() if f.type == "binary"]
        for field_name in binary_fields:
            if record[field_name]:
                try:
                    # Get the field's content
                    binary_content = record[field_name]
                    
                    # Try to determine mimetype based on filename or field name
                    mimetype = "application/octet-stream"
                    
                    # If there's a filename field, try to use it to determine mimetype
                    filename_field = f"{field_name}_filename"
                    if hasattr(record, filename_field) and record[filename_field]:
                        filename = record[filename_field]
                        if filename.endswith('.pdf'):
                            mimetype = "application/pdf"
                        elif filename.endswith(('.doc', '.docx')):
                            mimetype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        elif filename.endswith(('.ppt', '.pptx')):
                            mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        elif filename.endswith('.html'):
                            mimetype = "text/html"
                        elif filename.endswith(('.jpg', '.jpeg')):
                            mimetype = "image/jpeg"
                        elif filename.endswith('.png'):
                            mimetype = "image/png"
                    
                    results.append({
                        "field_name": field_name,
                        "mimetype": mimetype,
                        "rawcontent": binary_content
                    })
                    
                except Exception as e:
                    _logger.error(f"Error processing binary field {field_name}: {e}")
                    continue

        # Get attachments related to this record
        try:
            attachments = env["ir.attachment"].search([
                ("res_model", "=", record._name),
                ("res_id", "=", record.id)
            ])
            
            for attachment in attachments:
                # Skip attachments without data or filename
                if not attachment.datas or not attachment.name:
                    continue
                    
                # Get binary content
                binary_content = attachment.datas
                
                # Determine mimetype
                mimetype = attachment.mimetype or "application/octet-stream"
                
                results.append({
                    "field_name": f"Attachment: {attachment.name}",
                    "mimetype": mimetype,
                    "rawcontent": binary_content,
                    "attachment_id": attachment.id
                })
        except Exception as e:
            _logger.error(f"Error processing attachments: {e}")
            
        return results

    def _get_file_extension(self, filename):
        """Extract file extension from filename."""
        if not filename:
            return ""
        parts = filename.split('.')
        if len(parts) > 1:
            return parts[-1].lower()
        return ""

    def _parse_markdown(self, resource, field):
        """Parse Markdown content specifically."""
        field_name = field.get("field_name", "")
        raw_content = field.get("rawcontent", "")
        
        # If content is in bytes, decode it properly
        if isinstance(raw_content, bytes):
            content = self._safe_decode(raw_content)
        elif isinstance(raw_content, str) and raw_content.startswith("b'"):
            # This handles the case where bytes are represented as a string like b'content'
            try:
                # Remove the b'' wrapper and parse
                if raw_content.endswith("'"):
                    content_str = raw_content[2:-1]
                else:
                    content_str = raw_content[2:]  # Handle truncated strings
                # Handle escape sequences
                content = content_str.encode('utf-8').decode('unicode_escape')
            except Exception as e:
                self._log_error(f"Failed to decode markdown string representation: {str(e)}")
                content = raw_content
        else:
            content = raw_content
            
        # Ensure we have properly formatted Markdown
        if content and not content.strip().startswith("#"):
            # Add filename as a top-level header if the content doesn't start with a header
            return f"# {field_name}\n\n{content}"
        
        return content
